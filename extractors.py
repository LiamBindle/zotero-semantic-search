"""Text extraction from Zotero attachment files.

Each extractor returns a list of {"text": str, "location": str} dicts.
All imports are lazy so only the libraries needed for the file type in question
are loaded.  The single public entry point is `extract(path)`.
"""

import logging
import re
import zipfile
from pathlib import Path
from typing import Iterator

log = logging.getLogger(__name__)

CHUNK_CHARS = 2000
OVERLAP_CHARS = 200
MIN_CHUNK_CHARS = 80

SUPPORTED_SUFFIXES = frozenset({
    ".pdf",
    ".html", ".htm",
    ".docx", ".doc",
    ".odt",
    ".pptx", ".ppt",
    ".xlsx", ".xls",
    ".epub",
    ".txt", ".md",
    ".rtf",
})


# ── Shared utilities ───────────────────────────────────────────────────────────

def _ws(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


# DOIs appear only in reference list entries, never in body text.
_DOI_RE = re.compile(r'doi\.org/10\.|\bdoi:\s*10\.', re.IGNORECASE)


def _is_ref_block(text: str) -> bool:
    """Return True if this block is a bibliography entry (contains a DOI).

    Year-density heuristics are intentionally omitted: they have unacceptable
    false-positive risk for policy/regulation documents that list many dates,
    and the penalty for excluding a real paragraph is worse than including a
    reference entry.  DOI detection alone is essentially noise-free.
    """
    return bool(_DOI_RE.search(text))


def _sliding_window(text: str,
                    size: int = CHUNK_CHARS,
                    overlap: int = OVERLAP_CHARS) -> list[str]:
    text = _ws(text)
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end = min(start + size, len(text))
        if end < len(text):
            search_start = start + int(size * 0.8)
            boundary = max(
                text.rfind(". ", search_start, end),
                text.rfind(".\n", search_start, end),
            )
            if boundary != -1:
                end = boundary + 1
        chunk = text[start:end].strip()
        if len(chunk) >= MIN_CHUNK_CHARS:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - overlap
    return chunks


# ── PDF ────────────────────────────────────────────────────────────────────────

def _extract_pdf(path: str) -> list[dict]:
    import fitz  # pymupdf
    chunks: list[dict] = []
    doc = fitz.open(path)
    try:
        for page_num in range(len(doc)):
            location = f"Page {page_num + 1}"
            for block in doc[page_num].get_text("blocks"):
                if block[6] != 0:  # skip image blocks
                    continue
                text = _ws(block[4])
                if _is_ref_block(text):
                    continue
                for piece in _sliding_window(text):
                    chunks.append({"text": piece, "location": location})
    finally:
        doc.close()
    return chunks


# ── HTML ───────────────────────────────────────────────────────────────────────

_HTML_SKIP = frozenset({
    "style", "script", "noscript", "svg", "canvas", "template",
    "nav", "footer", "aside", "figure",
    "form", "button", "select", "option", "input", "textarea",
    "head", "meta", "link",
})
_HTML_HEADING = frozenset({"h1", "h2", "h3", "h4", "h5", "h6"})
# Leaf-block elements: emit their full text if they contain no nested block/heading children.
# div/section/article/main are included so that text held directly in a <div> (no <p> child)
# is not silently dropped — common in journal article HTML snapshots.
_HTML_BLOCK = frozenset({
    "p", "blockquote", "pre", "li", "dd", "dt", "td", "th", "caption",
    "div", "section", "article", "main",
})


def _walk_html(el) -> Iterator[tuple[bool, str]]:
    """Yield (is_heading, text) without double-counting nested elements."""
    tag = el.tag if isinstance(el.tag, str) else None
    if not tag or tag in _HTML_SKIP:
        return

    if tag in _HTML_HEADING:
        text = _ws(" ".join(el.itertext()))
        if text:
            yield (True, text)
        return

    if tag in _HTML_BLOCK:
        # Emit this element's text only if it contains no nested block/heading children,
        # preventing the same text from being yielded by both parent and child elements.
        has_block_child = any(
            c.tag in _HTML_BLOCK or c.tag in _HTML_HEADING
            for c in el if isinstance(c.tag, str)
        )
        if not has_block_child:
            text = _ws(" ".join(el.itertext()))
            if text:
                yield (False, text)
        else:
            for child in el:
                yield from _walk_html(child)
        return

    for child in el:
        yield from _walk_html(child)


def _chunks_from_html_bytes(data: bytes) -> list[dict]:
    from lxml import etree
    parser = etree.HTMLParser(remove_comments=True)
    root = etree.fromstring(data, parser)
    # Start from <body> to skip all <head> content (embedded fonts, base64 CSS, etc.)
    start = root.find(".//body")
    if start is None:
        start = root

    chunks: list[dict] = []
    current_heading = "Document"
    for is_heading, text in _walk_html(start):
        if is_heading:
            current_heading = text
            continue
        if _is_ref_block(text):
            continue
        for piece in _sliding_window(text):
            chunks.append({"text": piece, "location": current_heading})
    return chunks


def _extract_html(path: str) -> list[dict]:
    with open(path, "rb") as f:
        return _chunks_from_html_bytes(f.read())


# ── DOCX ───────────────────────────────────────────────────────────────────────

def _extract_docx(path: str) -> list[dict]:
    import docx
    doc = docx.Document(path)
    chunks: list[dict] = []
    current_heading = "Document"

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if "heading" in (para.style.name or "").lower():
            current_heading = text
        for piece in _sliding_window(text):
            chunks.append({"text": piece, "location": current_heading})

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    for piece in _sliding_window(text):
                        chunks.append({"text": piece, "location": current_heading})

    return chunks


# ── ODT ────────────────────────────────────────────────────────────────────────

def _extract_odt(path: str) -> list[dict]:
    from lxml import etree
    TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    H_TAG = f"{{{TEXT_NS}}}h"
    P_TAG = f"{{{TEXT_NS}}}p"

    with zipfile.ZipFile(path, "r") as zf:
        root = etree.fromstring(zf.read("content.xml"))

    chunks: list[dict] = []
    current_heading = "Document"
    for el in root.iter():
        tag = el.tag if isinstance(el.tag, str) else ""
        text = _ws("".join(el.itertext()))
        if not text:
            continue
        if tag == H_TAG:
            current_heading = text
            for piece in _sliding_window(text):
                chunks.append({"text": piece, "location": current_heading})
        elif tag == P_TAG:
            for piece in _sliding_window(text):
                chunks.append({"text": piece, "location": current_heading})
    return chunks


# ── PPTX ───────────────────────────────────────────────────────────────────────

def _extract_pptx(path: str) -> list[dict]:
    from pptx import Presentation
    chunks: list[dict] = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        location = f"Slide {i}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    for piece in _sliding_window(text):
                        chunks.append({"text": piece, "location": location})
    return chunks


# ── XLSX ───────────────────────────────────────────────────────────────────────

def _extract_xlsx(path: str) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    chunks: list[dict] = []
    try:
        for sheet_name in wb.sheetnames:
            location = f"Sheet: {sheet_name}"
            rows = []
            for row in wb[sheet_name].iter_rows():
                cells = [
                    str(c.value).strip()
                    for c in row
                    if c.value is not None and str(c.value).strip()
                ]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                for piece in _sliding_window("\n".join(rows)):
                    chunks.append({"text": piece, "location": location})
    finally:
        wb.close()
    return chunks


# ── EPUB ───────────────────────────────────────────────────────────────────────

def _extract_epub(path: str) -> list[dict]:
    import posixpath
    from lxml import etree

    OPF_NS = "http://www.idpf.org/2007/opf"
    OCF_NS = "urn:oasis:names:tc:opendocument:xmlns:container"

    chunks: list[dict] = []
    with zipfile.ZipFile(path, "r") as zf:
        # Locate the OPF package document via META-INF/container.xml
        c_root = etree.fromstring(zf.read("META-INF/container.xml"))
        rf = (c_root.find(f"{{{OCF_NS}}}rootfiles/{{{OCF_NS}}}rootfile")
              or c_root.find(".//rootfile"))
        if rf is None:
            return []

        opf_path = rf.get("full-path", "")
        opf_dir = posixpath.dirname(opf_path)
        opf_root = etree.fromstring(zf.read(opf_path))

        # Build id → href map for HTML items in the manifest
        manifest = {
            item.get("id"): item.get("href", "")
            for item in opf_root.findall(f"{{{OPF_NS}}}manifest/{{{OPF_NS}}}item")
            if "html" in item.get("media-type", "")
        }

        # Process spine items in reading order
        for itemref in opf_root.findall(f"{{{OPF_NS}}}spine/{{{OPF_NS}}}itemref"):
            href = manifest.get(itemref.get("idref"), "")
            if not href:
                continue
            full = posixpath.join(opf_dir, href).lstrip("/") if opf_dir else href
            try:
                chunks.extend(_chunks_from_html_bytes(zf.read(full)))
            except Exception as e:
                log.debug("EPUB chapter %s skipped: %s", full, e)

    return chunks


# ── Plain text / Markdown ──────────────────────────────────────────────────────

def _extract_text(path: str) -> list[dict]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return [{"text": p, "location": "Document"} for p in _sliding_window(text)]


# ── RTF ────────────────────────────────────────────────────────────────────────

def _extract_rtf(path: str) -> list[dict]:
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = rtf_to_text(f.read())
    return [{"text": p, "location": "Document"} for p in _sliding_window(text)]


# ── Dispatch ───────────────────────────────────────────────────────────────────

_DISPATCH: dict[str, object] = {
    ".pdf":  _extract_pdf,
    ".html": _extract_html,
    ".htm":  _extract_html,
    ".docx": _extract_docx,
    ".doc":  _extract_docx,
    ".odt":  _extract_odt,
    ".pptx": _extract_pptx,
    ".ppt":  _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,
    ".epub": _extract_epub,
    ".txt":  _extract_text,
    ".md":   _extract_text,
    ".rtf":  _extract_rtf,
}


def extract(path: str) -> list[dict]:
    """Extract text chunks from a file.

    Dispatches to the appropriate extractor based on file extension.
    Returns [] for unsupported formats or on unrecoverable errors.
    Each chunk is {"text": str, "location": str}.
    """
    suffix = Path(path).suffix.lower()
    fn = _DISPATCH.get(suffix)
    if fn is None:
        return []
    try:
        return fn(path)  # type: ignore[call-arg]
    except Exception as e:
        log.warning("Extraction failed for %s: %s", Path(path).name, e)
        return []
